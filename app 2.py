import streamlit as st
import os
import time
import json
import base64
from datetime import datetime
from openai import OpenAI
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import pandas as pd

# ---------------------------------------------------------
# إعدادات الصفحة
# ---------------------------------------------------------
st.set_page_config(page_title="سيرتي | Seeraty", layout="wide", page_icon="📄")

# CSS
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Tajawal:wght@400;500;700;800&display=swap');
html, body, [class*='css'] {font-family: 'Tajawal', sans-serif;}
.main-title {text-align: center; font-size: 48px; font-weight: 800; color: #0f172a; margin-bottom: 10px;}
.sub-title {text-align: center; font-size: 18px; color: #64748b; margin-bottom: 40px;}
.stButton>button {width: 100%; border-radius: 10px; height: 3.5em; background-color: #0f172a; color: white; font-weight: bold; border: none;}
.result-card {background-color: #ffffff; padding: 25px; border-radius: 12px; box-shadow: 0 4px 6px -1px rgba(0,0,0,0.1); border: 1px solid #e2e8f0; margin-bottom: 20px; color: #000000 !important;}
</style>
""", unsafe_allow_html=True)

# ---------------------------------------------------------
# سحب المفتاح من Secrets (الطريقة الآمنة)
# ---------------------------------------------------------
try:
    server_key = st.secrets["OPENAI_API_KEY"]
except:
    st.error("⚠️ خطأ: لم يتم إضافة المفتاح في إعدادات الاستضافة (Secrets).")
    st.stop()

# ---------------------------------------------------------
# الدوال المساعدة
# ---------------------------------------------------------
def extract_text(file):
    name = file.name.lower()
    text = ""
    img_b64 = None
    try:
        if name.endswith(".pdf"):
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages: text += page.extract_text() + "\n"
        elif name.endswith(".docx"):
            doc = Document(file)
            for p in doc.paragraphs: text += p.text + "\n"
        elif name.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        elif name.endswith(".txt"):
            text = file.read().decode("utf-8")
        elif name.endswith((".png", ".jpg", ".jpeg")):
            img_b64 = base64.b64encode(file.read()).decode("utf-8")
            text = "IMAGE_MODE"
        return text, img_b64
    except:
        return None, None

# ---------------------------------------------------------
# الواجهة
# ---------------------------------------------------------
st.markdown('<div class="main-title">سيرتي</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-title">نظام تحليل وتطوير السيرة الذاتية الاحترافي</div>', unsafe_allow_html=True)

uploaded_file = st.file_uploader("قم برفع السيرة الذاتية لبدء التحليل", type=["pdf", "docx", "pptx", "png", "jpg", "jpeg"])

if st.button("🚀 ابدأ التحليل الاحترافي"):
    if not uploaded_file:
        st.warning("⚠️ الرجاء رفع ملف أولاً.")
    else:
        status_box = st.status("جاري معالجة الملف...", expanded=True)
        try:
            status_box.write("📂 قراءة الملف...")
            content, img_data = extract_text(uploaded_file)
            
            if not content:
                status_box.update(label="فشل القراءة", state="error")
                st.error("الملف تالف أو لا يمكن قراءته.")
                st.stop()

            client = OpenAI(api_key=server_key)

            # الفرز
            status_box.write("🕵️‍♂️ التحقق من الملف...")
            check_msg = [{"role": "system", "content": "هل هذا الملف CV؟ رد بـ VALID أو INVALID"}]
            if content == "IMAGE_MODE":
                check_msg.append({"role": "user", "content": [{"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_data}"}}]})
            else:
                check_msg.append({"role": "user", "content": content[:2000]})
            
            check_res = client.chat.completions.create(model="gpt-4o", messages=check_msg)
            
            if "INVALID" in check_res.choices[0].message.content:
                status_box.update(label="مرفوض", state="error")
                st.error("🛑 عذراً، هذا الملف ليس سيرة ذاتية.")
                st.stop()

            # التحليل
            status_box.write("🧠 جاري التحليل العميق...")
            analyze_prompt = """
            تصرف كمستشار توظيف خبير. حلل السيرة الذاتية التالية واستخرج:
            1. توقع الراتب بالريال السعودي (نطاق).
            2. تقييم القوة (من 100).
            3. الأخطاء والنواقص الحقيقية.
            4. 3 شهادات مهنية مقترحة لرفع الراتب (بالاسم الإنجليزي).
            5. نصيحة ذهبية.
            نسق الرد بشكل جميل وجاهز للعرض.
            """
            
            # تجهيز رسالة التحليل
            analyze_msg = [{"role": "system", "content": analyze_prompt}]
            if content == "IMAGE_MODE":
                analyze_msg.append({"role": "user", "content": [{"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_data}"}}]})
            else:
                analyze_msg.append({"role": "user", "content": content[:4000]})

            final_res = client.chat.completions.create(model="gpt-4o", messages=analyze_msg)
            report = final_res.choices[0].message.content
            
            status_box.update(label="✅ تم التحليل!", state="complete", expanded=False)
            st.markdown("---")
            st.markdown(f'<div class="result-card">{report}</div>', unsafe_allow_html=True)

        except Exception as e:
            status_box.update(label="خطأ", state="error")
            st.error(f"حدث خطأ: {e}")
